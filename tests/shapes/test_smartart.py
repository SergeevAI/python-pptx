"""Unit-test suite for pptx.shapes.smartart module."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.parts.diagram import DiagramDataPart
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.smartart import SmartArt, SmartArtNode, _SmartArtNodes

from ..unitutil.cxml import element
from ..unitutil.mock import instance_mock


class DescribeSmartArt:
    """Unit-test suite for `pptx.shapes.smartart.SmartArt` object."""

    def it_can_be_constructed(self, data_part_):
        """Test SmartArt object construction."""
        smartart = SmartArt(data_part_, None)
        assert smartart._data_part is data_part_

    def it_provides_access_to_nodes(self, data_part_with_points_):
        """Test accessing nodes collection."""
        smartart = SmartArt(data_part_with_points_, None)
        nodes = smartart.nodes

        assert isinstance(nodes, _SmartArtNodes)

    def it_can_iterate_text_content(self, data_part_with_text_):
        """Test iterating text content."""
        smartart = SmartArt(data_part_with_text_, None)
        texts = list(smartart.iter_text())

        assert len(texts) > 0
        assert all(isinstance(text, str) for text in texts)

    def it_provides_text_content_as_list(self, data_part_with_text_):
        """Test getting text content as list."""
        smartart = SmartArt(data_part_with_text_, None)
        text_list = smartart.text_content

        assert isinstance(text_list, list)
        assert all(isinstance(text, str) for text in text_list)

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def data_part_(self, request):
        xml = (
            '<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
            "  <dgm:ptLst/>"
            "  <dgm:cxnLst/>"
            "</dgm:dataModel>"
        )
        data_model = parse_xml(xml)
        data_part = instance_mock(request, DiagramDataPart, data_model=data_model)
        return data_part

    @pytest.fixture
    def data_part_with_points_(self, request):
        xml = (
            '<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "  <dgm:ptLst>"
            '    <dgm:pt modelId="{id1}"/>'
            '    <dgm:pt modelId="{id2}"/>'
            "  </dgm:ptLst>"
            "  <dgm:cxnLst/>"
            "</dgm:dataModel>"
        )
        data_model = parse_xml(xml)
        data_part = instance_mock(request, DiagramDataPart, data_model=data_model)
        return data_part

    @pytest.fixture
    def data_part_with_text_(self, request):
        xml = (
            '<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "  <dgm:ptLst>"
            '    <dgm:pt modelId="{id1}">'
            "      <dgm:t>"
            "        <a:p><a:r><a:t>test1</a:t></a:r></a:p>"
            "      </dgm:t>"
            "    </dgm:pt>"
            '    <dgm:pt modelId="{id2}">'
            "      <dgm:t>"
            "        <a:p><a:r><a:t>test2</a:t></a:r></a:p>"
            "      </dgm:t>"
            "    </dgm:pt>"
            "  </dgm:ptLst>"
            "  <dgm:cxnLst/>"
            "</dgm:dataModel>"
        )
        data_model = parse_xml(xml)
        data_part = instance_mock(request, DiagramDataPart, data_model=data_model)
        return data_part


class Describe_SmartArtNodes:
    """Unit-test suite for `pptx.shapes.smartart._SmartArtNodes` object."""

    def it_can_iterate_nodes(self):
        """Test iterating over nodes."""
        xml = (
            '<dgm:ptLst xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '  <dgm:pt modelId="{id1}"/>'
            '  <dgm:pt modelId="{id2}"/>'
            '  <dgm:pt modelId="{id3}" type="pres"/>'  # Should be skipped
            "  </dgm:ptLst>"
        )
        pt_lst = parse_xml(xml)
        nodes = _SmartArtNodes(pt_lst, None, None)
        node_list = list(nodes)

        # Only 2 nodes (pres type is skipped)
        assert len(node_list) == 2
        assert all(isinstance(node, SmartArtNode) for node in node_list)

    def it_can_count_nodes(self):
        """Test counting nodes."""
        xml = (
            '<dgm:ptLst xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '  <dgm:pt modelId="{id1}"/>'
            '  <dgm:pt modelId="{id2}"/>'
            '  <dgm:pt modelId="{id3}" type="pres"/>'
            "  </dgm:ptLst>"
        )
        pt_lst = parse_xml(xml)
        nodes = _SmartArtNodes(pt_lst, None, None)

        assert len(nodes) == 2

    def it_filters_out_presentation_nodes(self):
        """Test that presentation, transition nodes are filtered."""
        xml = (
            '<dgm:ptLst xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '  <dgm:pt modelId="{id1}"/>'  # Include
            '  <dgm:pt modelId="{id2}" type="pres"/>'  # Exclude
            '  <dgm:pt modelId="{id3}" type="parTrans"/>'  # Exclude
            '  <dgm:pt modelId="{id4}" type="sibTrans"/>'  # Exclude
            '  <dgm:pt modelId="{id5}" type="doc"/>'  # Include
            "  </dgm:ptLst>"
        )
        pt_lst = parse_xml(xml)
        nodes = _SmartArtNodes(pt_lst, None, None)
        node_list = list(nodes)

        assert len(node_list) == 2
        model_ids = [node.model_id for node in node_list]
        assert "{id1}" in model_ids
        assert "{id5}" in model_ids


class DescribeSmartArtNode:
    """Unit-test suite for `pptx.shapes.smartart.SmartArtNode` object."""

    def it_provides_text_property(self):
        """Test accessing node text."""
        xml = (
            '<dgm:pt xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' modelId="{test}">'
            "  <dgm:t>"
            "    <a:p><a:r><a:t>Node Text</a:t></a:r></a:p>"
            "  </dgm:t>"
            "</dgm:pt>"
        )
        pt = parse_xml(xml)
        node = SmartArtNode(pt, None, None)

        assert node.text == "Node Text"

    def it_provides_model_id_property(self):
        """Test accessing node model ID."""
        xml = (
            '<dgm:pt xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' modelId="{TEST-MODEL-ID}"/>'
        )
        pt = parse_xml(xml)
        node = SmartArtNode(pt, None, None)

        assert node.model_id == "{TEST-MODEL-ID}"

    def it_can_set_text_property(self):
        """Test setting node text."""
        xml = (
            '<dgm:pt xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' modelId="{test}">'
            "  <dgm:t>"
            "    <a:p><a:r><a:t>Old Text</a:t></a:r></a:p>"
            "  </dgm:t>"
            "</dgm:pt>"
        )
        pt = parse_xml(xml)
        node = SmartArtNode(pt, None, None)

        node.text = "New Text"
        assert node.text == "New Text"

    def it_provides_node_type_property(self):
        """Test accessing node type."""
        xml = (
            '<dgm:pt xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' type="doc"/>'
        )
        pt = parse_xml(xml)
        node = SmartArtNode(pt, None, None)

        assert node.node_type == "doc"


class DescribeGraphicFrameSmartArtIntegration:
    """Integration tests for GraphicFrame with SmartArt."""

    def it_knows_when_it_contains_smartart(self):
        """Test detecting SmartArt in graphic frame."""
        xml = (
            '<p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"'
            ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            "  <p:nvGraphicFramePr>"
            '    <p:cNvPr id="1" name="SmartArt"/>'
            "    <p:cNvGraphicFramePr/>"
            "    <p:nvPr/>"
            "  </p:nvGraphicFramePr>"
            "  <p:xfrm>"
            '    <a:off x="0" y="0"/>'
            '    <a:ext cx="100" cy="100"/>'
            "  </p:xfrm>"
            "  <a:graphic>"
            '    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
            '      <dgm:relIds r:dm="rId1" r:lo="rId2" r:qs="rId3" r:cs="rId4"/>'
            "    </a:graphicData>"
            "  </a:graphic>"
            "</p:graphicFrame>"
        )
        graphic_frame_elm = parse_xml(xml)
        graphic_frame = GraphicFrame(graphic_frame_elm, None)

        assert graphic_frame.has_smartart is True
        assert graphic_frame.has_chart is False
        assert graphic_frame.has_table is False

    def it_reports_correct_shape_type_for_smartart(self):
        """Test that SmartArt graphic frame reports correct shape type."""
        xml = (
            '<p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
            "  <p:nvGraphicFramePr>"
            '    <p:cNvPr id="1" name="SmartArt"/>'
            "    <p:cNvGraphicFramePr/>"
            "    <p:nvPr/>"
            "  </p:nvGraphicFramePr>"
            "  <p:xfrm>"
            '    <a:off x="0" y="0"/>'
            '    <a:ext cx="100" cy="100"/>'
            "  </p:xfrm>"
            "  <a:graphic>"
            '    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>'
            "  </a:graphic>"
            "</p:graphicFrame>"
        )
        graphic_frame_elm = parse_xml(xml)
        graphic_frame = GraphicFrame(graphic_frame_elm, None)

        assert graphic_frame.shape_type == MSO_SHAPE_TYPE.SMART_ART


class DescribeSmartArtAcceptanceTest:
    """Acceptance test using real PPTX file."""

    def it_can_detect_smartart_in_test_file(self):
        """Test detecting SmartArt in smartart-test.pptx."""
        try:
            prs = Presentation("tests/test_files/smartart-test.pptx")
        except FileNotFoundError:
            pytest.skip("smartart-test.pptx not found")

        # Slide 2 (index 1) contains SmartArt
        slide = prs.slides[1]
        shape = slide.shapes[0]

        assert isinstance(shape, GraphicFrame)
        assert shape.has_smartart is True

    def it_can_extract_text_from_smartart(self):
        """Test extracting text content from SmartArt."""
        try:
            prs = Presentation("tests/test_files/smartart-test.pptx")
        except FileNotFoundError:
            pytest.skip("smartart-test.pptx not found")

        slide = prs.slides[1]
        shape = slide.shapes[0]

        if not shape.has_smartart:
            pytest.skip("Shape is not SmartArt")

        smartart = shape.smartart
        texts = smartart.text_content

        # Should find 4 "test" strings
        test_texts = [t for t in texts if t == "test"]
        assert len(test_texts) == 4

    def it_can_iterate_smartart_nodes(self):
        """Test iterating over SmartArt nodes."""
        try:
            prs = Presentation("tests/test_files/smartart-test.pptx")
        except FileNotFoundError:
            pytest.skip("smartart-test.pptx not found")

        slide = prs.slides[1]
        shape = slide.shapes[0]

        if not shape.has_smartart:
            pytest.skip("Shape is not SmartArt")

        smartart = shape.smartart
        nodes = list(smartart.nodes)

        assert len(nodes) > 0
        for node in nodes:
            assert isinstance(node, SmartArtNode)
            assert node.model_id is not None
